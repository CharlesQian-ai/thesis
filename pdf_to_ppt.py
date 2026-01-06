import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def convert_pdf_to_ppt(pdf_path, ppt_path):
    print(f"Opening PDF: {pdf_path}")
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"Error opening PDF: {e}")
        return

    # 创建 PPT 对象
    prs = Presentation()
    
    # 设置 PPT 尺寸为宽屏 16:9 (默认是 4:3)
    # 1280x720 像素对应 16:9 比例
    # PowerPoint 默认宽屏尺寸通常是 13.333 inches x 7.5 inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    print(f"Total pages: {len(doc)}")

    for page_num in range(len(doc)):
        # print(f"Processing page {page_num + 1}...")
        page = doc.load_page(page_num)
        
        # 将 PDF 页面转换为图片 (pixmap)
        # matrix=fitz.Matrix(3, 3) 表示放大 3 倍，提高清晰度
        pix = page.get_pixmap(matrix=fitz.Matrix(3, 3))
        
        image_path = f"temp_slide_{page_num}.png"
        pix.save(image_path)
        
        # 创建空白幻灯片 (layout 6 是空白页)
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 将图片填满整个幻灯片
        slide.shapes.add_picture(
            image_path, 
            0, 
            0, 
            width=prs.slide_width, 
            height=prs.slide_height
        )
        
        # 删除临时图片
        if os.path.exists(image_path):
            os.remove(image_path)

    # 确保输出目录存在
    os.makedirs(os.path.dirname(ppt_path), exist_ok=True)
    
    # 保存 PPT
    prs.save(ppt_path)
    print(f"Successfully saved PPT to: {ppt_path}")

if __name__ == "__main__":
    # 默认测试路径
    PDF_PATH = r"d:\Projects\开题答辩PPT\proposal\docs\开题答辩PPT.pdf"
    OUTPUT_PPT_PATH = r"d:\Projects\开题答辩PPT\proposal\docs\开题答辩PPT.pptx"
    
    if os.path.exists(PDF_PATH):
        convert_pdf_to_ppt(PDF_PATH, OUTPUT_PPT_PATH)
    else:
        print(f"PDF not found: {PDF_PATH}")
